# Copyright 2025-2026 Hanlian Lu. SPDX-License-Identifier: Apache-2.0
"""Tests for deterministic binary resource conversion and OOXML preflight."""

from __future__ import annotations

import asyncio
import io
import socket
import tomllib
import zipfile
from pathlib import Path

import openpyxl
import pytest
from docx import Document
from openpyxl.drawing.image import Image as XLImage
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from dlightrag.engine.answer.resources.converters import (
    ConvertedResource,
    ResourceConversionError,
    UnsafeArchiveError,
    convert_resource,
    is_convertible,
)
from dlightrag.engine.answer.resources.converters import (
    _validate_archive_sizes as validate_archive_sizes,
)
from dlightrag.engine.answer.resources.models import ResourceInput
from dlightrag.engine.answer.resources.registry import ResourceRegistry

_FIXTURES = Path(__file__).parent / "fixtures" / "resources"

_MIB = 1024 * 1024


def _png(color: tuple[int, int, int]) -> bytes:
    buffer = io.BytesIO()
    Image.new("RGB", (8, 8), color).save(buffer, "PNG")
    return buffer.getvalue()


def _docx_bytes(*, heading: str, body: str, image: bytes | None = None) -> bytes:
    document = Document()
    document.add_heading(heading, level=1)
    document.add_paragraph(body)
    if image is not None:
        document.add_picture(io.BytesIO(image))
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _pptx_bytes(*, title: str, image: bytes | None = None) -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    title_shape = slide.shapes.title
    assert title_shape is not None
    title_shape.text = title
    if image is not None:
        slide.shapes.add_picture(io.BytesIO(image), Inches(1), Inches(1))
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _xlsx_bytes(
    *,
    sheet_title: str,
    cells: dict[str, str | int | float],
    image: bytes | None = None,
    image_cell: str = "B2",
) -> bytes:
    workbook = openpyxl.Workbook()
    worksheet = workbook.active
    assert worksheet is not None
    worksheet.title = sheet_title
    for coordinate, value in cells.items():
        worksheet[coordinate] = value
    if image is not None:
        worksheet.add_image(XLImage(io.BytesIO(image)), image_cell)
    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"


async def test_html_conversion_yields_semantic_text() -> None:
    content = (_FIXTURES / "sample.html").read_bytes()

    result = await convert_resource(content, filename="sample.html", declared_mime="text/html")

    assert isinstance(result, ConvertedResource)
    assert "Acme Quarterly Summary" in result.text
    assert "Revenue" in result.text
    assert "<h1>" not in result.text
    assert result.visuals == ()


async def test_csv_conversion_yields_table_rows() -> None:
    content = (_FIXTURES / "sample.csv").read_bytes()

    result = await convert_resource(content, filename="sample.csv", declared_mime="text/csv")

    assert "Widget" in result.text
    assert "1200" in result.text


async def test_pdf_conversion_extracts_text() -> None:
    content = (_FIXTURES / "sample.pdf").read_bytes()

    result = await convert_resource(content, filename="report.pdf", declared_mime="application/pdf")

    assert "Quarterly Revenue Report" in result.text
    assert "subscription" in result.text


async def test_docx_image_becomes_visual_handle_without_base64() -> None:
    content = _docx_bytes(
        heading="Master Agreement",
        body="Total liability is capped at fees paid.",
        image=_png((220, 20, 20)),
    )

    result = await convert_resource(content, filename="contract.docx", declared_mime=DOCX_MIME)

    assert "Total liability" in result.text
    assert "data:image" not in result.text
    assert "base64" not in result.text
    assert len(result.visuals) == 1
    visual = result.visuals[0]
    assert visual.media_type == "image/png"
    assert visual.data.startswith(b"\x89PNG")
    assert visual.handle_id in result.text


async def test_pptx_image_becomes_visual_handle() -> None:
    content = _pptx_bytes(title="Roadmap Milestones", image=_png((0, 180, 0)))

    result = await convert_resource(content, filename="deck.pptx", declared_mime=PPTX_MIME)

    assert "Roadmap Milestones" in result.text
    assert "data:image" not in result.text
    assert len(result.visuals) == 1
    assert result.visuals[0].data.startswith(b"\x89PNG")


async def test_xlsx_image_carries_sheet_and_cell_anchor() -> None:
    content = _xlsx_bytes(
        sheet_title="Financials",
        cells={"A1": "Revenue", "B1": 1200},
        image=_png((10, 10, 220)),
        image_cell="B2",
    )

    result = await convert_resource(content, filename="model.xlsx", declared_mime=XLSX_MIME)

    assert "Revenue" in result.text
    assert "data:image" not in result.text
    assert len(result.visuals) == 1
    visual = result.visuals[0]
    assert visual.anchor == "Financials!B2"
    assert visual.data.startswith(b"\x89PNG")


async def test_conversion_performs_no_network_fetch(monkeypatch: pytest.MonkeyPatch) -> None:
    def _blocked_connect(*args: object, **kwargs: object) -> None:
        raise AssertionError("conversion attempted a network connection")

    monkeypatch.setattr(socket.socket, "connect", _blocked_connect)
    html = b'<html><body><p>See <img src="https://evil.example/x.png"></p></body></html>'

    result = await convert_resource(html, filename="page.html", declared_mime="text/html")

    assert isinstance(result, ConvertedResource)


def test_is_convertible_routing() -> None:
    assert is_convertible("a.docx", None)
    assert is_convertible(None, XLSX_MIME)
    assert is_convertible("data.CSV", None)
    assert not is_convertible("notes.txt", "text/plain")
    assert not is_convertible(None, "application/json")


async def test_unconvertible_resource_is_rejected() -> None:
    with pytest.raises(ResourceConversionError):
        await convert_resource(b"plain", filename="notes.txt", declared_mime="text/plain")


def test_archive_entry_count_limit() -> None:
    sizes = [(0, 0)] * 10_001

    with pytest.raises(UnsafeArchiveError):
        validate_archive_sizes(sizes)


def test_archive_per_entry_size_limit() -> None:
    sizes = [(1024, 1024), (200 * _MIB, 4096)]

    with pytest.raises(UnsafeArchiveError):
        validate_archive_sizes(sizes)


def test_archive_total_size_limit() -> None:
    sizes = [(90 * _MIB, 90 * _MIB) for _ in range(6)]

    with pytest.raises(UnsafeArchiveError):
        validate_archive_sizes(sizes)


def test_archive_expansion_ratio_limit() -> None:
    sizes = [(50 * _MIB, 1024)]

    with pytest.raises(UnsafeArchiveError):
        validate_archive_sizes(sizes)


def test_archive_within_limits_is_accepted() -> None:
    validate_archive_sizes([(1024, 512), (2048, 1024)])


async def test_real_ratio_bomb_docx_is_rejected() -> None:
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", b"\x00" * (2 * _MIB))
    bomb = buffer.getvalue()

    with pytest.raises(UnsafeArchiveError):
        await convert_resource(bomb, filename="bomb.docx", declared_mime=DOCX_MIME)


async def test_eight_concurrent_conversions_do_not_cross_contaminate() -> None:
    payloads = [
        _docx_bytes(
            heading=f"Report {index}",
            body=f"Marker token unique-{index}.",
            image=_png((index * 20 % 256, 30, 200)),
        )
        for index in range(8)
    ]

    results = await asyncio.gather(
        *(
            convert_resource(payload, filename=f"r{index}.docx", declared_mime=DOCX_MIME)
            for index, payload in enumerate(payloads)
        )
    )

    for index, result in enumerate(results):
        assert f"unique-{index}" in result.text
        other_markers = [f"unique-{other}" for other in range(8) if other != index]
        assert not any(marker in result.text for marker in other_markers)
        assert len(result.visuals) == 1


def test_markitdown_extras_are_selected_only() -> None:
    dependencies = tomllib.loads(Path("pyproject.toml").read_text(encoding="utf-8"))["project"][
        "dependencies"
    ]
    markitdown = [dep for dep in dependencies if dep.startswith("markitdown")]

    assert len(markitdown) == 1
    extras = set(markitdown[0].split("[", 1)[1].split("]", 1)[0].split(","))
    assert extras == {"pdf", "docx", "pptx", "xlsx"}
    assert "all" not in markitdown[0]
    assert any(dep.startswith("pypdfium2") for dep in dependencies)
    assert any(dep.startswith("openpyxl") for dep in dependencies)


def test_no_forbidden_binary_reader_extras_in_lock() -> None:
    lock = Path("uv.lock").read_text(encoding="utf-8")

    forbidden = [
        "azure-ai-documentintelligence",
        "easyocr",
        "pytesseract",
        "pydub",
        "speechrecognition",
        "youtube-transcript-api",
        "xlrd",
    ]
    present = [name for name in forbidden if f'name = "{name}"' in lock]
    assert present == []


async def test_registry_read_routes_docx_to_visual_handle() -> None:
    content = _docx_bytes(
        heading="Master Agreement",
        body="Total liability is capped.",
        image=_png((200, 10, 10)),
    )
    async with ResourceRegistry() as registry:
        resource_id = registry.register(
            ResourceInput(filename="contract.docx", content=content, declared_mime=DOCX_MIME)
        )
        result = await registry.read(resource_id, max_window_tokens=1_000)

    assert "Total liability" in result.content
    assert "data:image" not in result.content
    assert len(result.visual_handles) == 1
    assert result.visual_handles[0].handle_id in result.content


async def test_registry_read_surfaces_xlsx_cell_anchor_label() -> None:
    content = _xlsx_bytes(
        sheet_title="Financials",
        cells={"A1": "Revenue", "B1": 1200},
        image=_png((10, 10, 220)),
        image_cell="B2",
    )
    async with ResourceRegistry() as registry:
        resource_id = registry.register(
            ResourceInput(filename="model.xlsx", content=content, declared_mime=XLSX_MIME)
        )
        result = await registry.read(resource_id, max_window_tokens=1_000)

    assert len(result.visual_handles) == 1
    assert result.visual_handles[0].label == "Financials!B2"


async def test_registry_text_resource_has_no_visual_handles() -> None:
    async with ResourceRegistry() as registry:
        resource_id = registry.register(
            ResourceInput(filename="notes.txt", content=b"plain notes", declared_mime="text/plain")
        )
        result = await registry.read(resource_id, max_window_tokens=1_000)

    assert result.content == "plain notes"
    assert result.visual_handles == ()
